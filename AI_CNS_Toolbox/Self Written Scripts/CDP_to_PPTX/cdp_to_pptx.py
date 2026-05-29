#!/usr/bin/env python3
"""
cdp_to_pptx.py

Parse the output of `show cdp neighbors detail` (Cisco IOS / NX-OS) and
generate a PowerPoint file containing:
    Slide 1: A topology diagram (local device in the center, neighbors
             around it) using Cisco icons.
    Slide 2: A connectivity matrix with all the important information.

This module is callable both from the command-line and from the Flask
web GUI via :func:`generate_pptx_from_cdp`.
"""

from __future__ import annotations

import argparse
import logging
import math
import os
import re
import sys
from dataclasses import dataclass, field
from typing import List, Optional

log = logging.getLogger("cdp_to_pptx")

# Make sure cairosvg can find the Homebrew-installed Cairo library on macOS
# without the user having to set DYLD_FALLBACK_LIBRARY_PATH manually.
for _candidate in ("/opt/homebrew/lib", "/usr/local/lib"):
    if os.path.isdir(_candidate):
        _current = os.environ.get("DYLD_FALLBACK_LIBRARY_PATH", "")
        if _candidate not in _current.split(":"):
            os.environ["DYLD_FALLBACK_LIBRARY_PATH"] = (
                f"{_candidate}:{_current}" if _current else _candidate
            )

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
CACHE_DIR = os.path.join(SCRIPT_DIR, ".icon_cache")


def _resolve_icon_dir() -> str:
    """Find the icon folder.

    Searches (in order):
      1. ``$CDP_ICON_DIR`` environment variable.
      2. ``<script_dir>/<candidate-name>``
      3. Walks up to 5 parent directories looking for a sibling folder with
         one of the candidate names.

    Candidate names handle both the original ``Cisco ICONS`` (with a space)
    and the underscore variant ``Cisco_ICONS`` that ships in this repo.
    """
    candidates = ("Cisco_ICONS", "Cisco ICONS", "cisco_icons", "cisco icons")

    env = os.environ.get("CDP_ICON_DIR")
    if env and os.path.isdir(env):
        return env

    for name in candidates:
        local = os.path.join(SCRIPT_DIR, name)
        if os.path.isdir(local):
            return local

    here = SCRIPT_DIR
    for _ in range(6):
        here = os.path.dirname(here)
        if not here or here == os.path.dirname(here):
            break
        for name in candidates:
            candidate = os.path.join(here, name)
            if os.path.isdir(candidate):
                return candidate

    # Fallback (may not exist; handled gracefully downstream).
    return os.path.join(SCRIPT_DIR, candidates[0])


ICON_DIR = _resolve_icon_dir()

ICON_FILES = {
    "router":   "router.svg",
    "switch":   "switch.svg",
    "l3switch": "switch.svg",
    "wlc":      "accesspoint.svg",
    "firewall": "firewall.svg",
    "host":     "accesspoint.svg",
    "unknown":  "switch.svg",
}

TYPE_LABELS = {
    "router":   "Router",
    "switch":   "Switch",
    "l3switch": "L3 Switch",
    "wlc":      "WLC",
    "firewall": "Firewall",
    "host":     "Host",
    "unknown":  "",
}

LOGO_SVG = "cisco_logo_icon_169399.svg"

# Color palette
CISCO_BLUE  = RGBColor(0x00, 0x51, 0x9E)
CISCO_DARK  = RGBColor(0x00, 0x2E, 0x5C)
LINK_COLOR  = RGBColor(0x4A, 0x6F, 0xA5)
TEXT_DARK   = RGBColor(0x1E, 0x2A, 0x3A)
TEXT_MUTED  = RGBColor(0x55, 0x66, 0x77)
BG_LIGHT    = RGBColor(0xF4, 0xF7, 0xFB)
BG_ALT      = RGBColor(0xE8, 0xEF, 0xF7)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

# Fallback colors for device icons when SVG/Cairo are unavailable.
DEVICE_FALLBACK_COLORS = {
    "router":   RGBColor(0x1F, 0x77, 0xB4),
    "switch":   RGBColor(0x2C, 0xA0, 0x2C),
    "l3switch": RGBColor(0x17, 0x6F, 0x50),
    "wlc":      RGBColor(0x9B, 0x59, 0xB6),
    "firewall": RGBColor(0xD9, 0x53, 0x4F),
    "host":     RGBColor(0x7F, 0x7F, 0x7F),
    "unknown":  RGBColor(0x55, 0x66, 0x77),
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ICON_SIZE = Inches(1.15)
LABEL_W   = Inches(2.6)
LABEL_H   = Inches(0.6)


# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------

@dataclass
class Neighbor:
    device_id: str = ""
    system_name: str = ""
    platform: str = ""
    capabilities: str = ""
    local_interface: str = ""
    remote_interface: str = ""
    ipv4_addresses: List[str] = field(default_factory=list)
    ipv6_addresses: List[str] = field(default_factory=list)
    mgmt_ipv4: List[str] = field(default_factory=list)
    mgmt_ipv6: List[str] = field(default_factory=list)
    version: str = ""
    native_vlan: str = ""
    duplex: str = ""

    @property
    def short_name(self) -> str:
        name = self.system_name or self.device_id or "unknown"
        name = re.sub(r"\s*\(.*?\)\s*$", "", name)
        name = name.split(".")[0]
        return name

    @property
    def primary_ipv4(self) -> str:
        if self.mgmt_ipv4:
            return self.mgmt_ipv4[0]
        if self.ipv4_addresses:
            return self.ipv4_addresses[0]
        return ""

    @property
    def device_type(self) -> str:
        return classify_device(self.platform, self.capabilities)


# ---------------------------------------------------------------------------
# Parsing
# ---------------------------------------------------------------------------

ENTRY_SEP = re.compile(r"^-{3,}\s*$", re.MULTILINE)


def _grab(pattern: str, text: str, flags: int = re.IGNORECASE) -> Optional[str]:
    m = re.search(pattern, text, flags)
    return m.group(1).strip() if m else None


def parse_cdp_output(text: str) -> List[Neighbor]:
    parts = ENTRY_SEP.split(text)
    neighbors: List[Neighbor] = []
    for chunk in parts:
        if "Device ID" not in chunk:
            continue
        n = _parse_entry(chunk)
        if n.device_id:
            neighbors.append(n)
    return neighbors


def _parse_entry(chunk: str) -> Neighbor:
    n = Neighbor()
    n.device_id     = _grab(r"Device ID\s*:?\s*(.+)", chunk) or ""
    n.system_name   = _grab(r"System Name\s*:?\s*(.+)", chunk) or ""
    n.platform      = _grab(r"Platform\s*:?\s*([^,]+),", chunk) or ""
    n.platform      = re.sub(r"^(cisco|Cisco)\s+", "", n.platform).strip()
    n.capabilities  = _grab(r"Capabilities\s*:?\s*(.+)", chunk) or ""
    n.local_interface  = _grab(r"^Interface\s*:\s*([^,]+),", chunk, re.IGNORECASE | re.MULTILINE) or ""
    n.remote_interface = _grab(r"Port ID \(outgoing port\)\s*:?\s*(.+)", chunk) or ""
    n.native_vlan   = _grab(r"Native VLAN\s*:?\s*(\S+)", chunk) or ""
    n.duplex        = _grab(r"Duplex\s*:?\s*(\S+)", chunk) or ""

    vmatch = re.search(r"Version\s*:\s*\n\s*(.+)", chunk)
    if vmatch:
        n.version = vmatch.group(1).strip()

    iface_block = _section(chunk, r"Interface address\(es\)", stop=r"(Platform|Mgmt address\(es\)|Local Interface MAC)")
    mgmt_block  = _section(chunk, r"Mgmt address\(es\)",      stop=r"(Local Interface MAC|DC_UUID|$)")

    n.ipv4_addresses = re.findall(r"IPv4 Address\s*:\s*([0-9.]+)", iface_block)
    n.ipv6_addresses = re.findall(r"IPv6 Address\s*:\s*([0-9a-fA-F:]+)", iface_block)
    n.mgmt_ipv4      = re.findall(r"IPv4 Address\s*:\s*([0-9.]+)", mgmt_block)
    n.mgmt_ipv6      = re.findall(r"IPv6 Address\s*:\s*([0-9a-fA-F:]+)", mgmt_block)

    return n


def _section(chunk: str, start_re: str, stop: str) -> str:
    m = re.search(start_re, chunk, re.IGNORECASE)
    if not m:
        return ""
    rest = chunk[m.end():]
    s = re.search(stop, rest)
    return rest[:s.start()] if s else rest


# ---------------------------------------------------------------------------
# Device classification
# ---------------------------------------------------------------------------

def classify_device(platform: str, capabilities: str) -> str:
    p = (platform or "").upper()
    c = (capabilities or "").upper()

    if any(t in p for t in ("ASA", "FTD", "FPR", "FIREPOWER", "FORTI", "PALO")):
        return "firewall"
    if any(t in p for t in ("AIR-CT", "C9800", "WLC", "AIR-WLC")):
        return "wlc"
    if p.startswith("AIR-") or "ACCESS POINT" in c:
        return "host"

    is_router = "ROUTER" in c
    is_switch = "SWITCH" in c

    if is_switch and is_router:
        return "l3switch"
    if is_switch:
        return "switch"
    if is_router:
        return "router"
    if "HOST" in c:
        return "host"
    return "unknown"


# ---------------------------------------------------------------------------
# Icon handling (with fallbacks if SVG icons or Cairo are unavailable)
# ---------------------------------------------------------------------------

# Warnings collected during the most recent call to generate_pptx_from_cdp().
_ICON_WARNINGS: List[str] = []
_RENDERER_CACHE: dict = {"backend": None, "render": None, "error": None}


def _add_warning(msg: str) -> None:
    if msg not in _ICON_WARNINGS:
        _ICON_WARNINGS.append(msg)
    log.warning("[cdp_to_pptx] %s", msg)


def _get_renderer():
    """Return (backend_name, render_callable) for SVG->PNG conversion.

    Tries ``cairosvg`` first, then falls back to ``svglib`` + ``reportlab``
    which is pure-Python and does not require system Cairo (handy on a
    fresh Linux box without ``libcairo2``).
    """
    if _RENDERER_CACHE["backend"] is not None:
        return _RENDERER_CACHE["backend"], _RENDERER_CACHE["render"]
    if _RENDERER_CACHE["error"] is not None:
        return None, None

    # Try cairosvg.
    try:
        import cairosvg  # type: ignore

        def _render_cairo(svg_path: str, png_path: str, size: int) -> None:
            cairosvg.svg2png(
                url=svg_path, write_to=png_path,
                output_width=size, output_height=size,
            )

        _RENDERER_CACHE["backend"] = "cairosvg"
        _RENDERER_CACHE["render"] = _render_cairo
        log.info("[cdp_to_pptx] SVG renderer: cairosvg")
        return "cairosvg", _render_cairo
    except Exception as e:
        _add_warning(
            "cairosvg unavailable: " + repr(e)
            + ". On Ubuntu install the system library with "
            + "`sudo apt-get install -y libcairo2 libpango-1.0-0 libpangocairo-1.0-0`, "
            + "or rely on the svglib fallback (`pip install svglib reportlab`)."
        )

    # Try svglib + reportlab (+ Pillow).
    try:
        from svglib.svglib import svg2rlg  # type: ignore
        from reportlab.graphics import renderPM  # type: ignore

        def _render_svglib(svg_path: str, png_path: str, size: int) -> None:
            drawing = svg2rlg(svg_path)
            if drawing is None:
                raise RuntimeError(f"svglib could not parse {svg_path}")
            w, h = float(drawing.width or size), float(drawing.height or size)
            scale = size / max(w, h) if max(w, h) > 0 else 1.0
            drawing.width = w * scale
            drawing.height = h * scale
            drawing.scale(scale, scale)
            renderPM.drawToFile(drawing, png_path, fmt="PNG")

        _RENDERER_CACHE["backend"] = "svglib"
        _RENDERER_CACHE["render"] = _render_svglib
        log.info("[cdp_to_pptx] SVG renderer: svglib + reportlab")
        return "svglib", _render_svglib
    except Exception as e:
        _add_warning(
            "svglib fallback unavailable: " + repr(e)
            + ". Install with `pip install svglib reportlab pillow`."
        )

    _RENDERER_CACHE["error"] = "no_renderer"
    _add_warning(
        "No SVG renderer available; falling back to colored rectangles "
        "for device icons."
    )
    return None, None


def _ensure_png(svg_name: str, size: int = 512) -> Optional[str]:
    svg_path = os.path.join(ICON_DIR, svg_name)
    if not os.path.exists(svg_path):
        _add_warning(f"Icon file not found: {svg_path}")
        return None

    backend, render = _get_renderer()
    if backend is None:
        return None

    try:
        os.makedirs(CACHE_DIR, exist_ok=True)
    except OSError as e:
        _add_warning(f"Cannot create icon cache dir {CACHE_DIR}: {e}")
        return None

    base = os.path.splitext(os.path.basename(svg_name))[0]
    png_path = os.path.join(CACHE_DIR, f"{base}_{backend}_{size}.png")

    needs_render = (
        not os.path.exists(png_path)
        or os.path.getmtime(png_path) < os.path.getmtime(svg_path)
    )
    if needs_render:
        try:
            render(svg_path, png_path, size)
        except Exception as e:
            _add_warning(
                f"Failed to render {svg_name} with {backend}: {e!r}"
            )
            return None
    return png_path


def _icon_png(dev_type: str) -> Optional[str]:
    fname = ICON_FILES.get(dev_type, ICON_FILES["unknown"])
    return _ensure_png(fname)


def _add_text(slide, left, top, width, height, text,
              size=10, bold=False, align=PP_ALIGN.CENTER,
              color=TEXT_DARK, fill=None, line_color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(40000)
    tf.margin_top = tf.margin_bottom = Emu(20000)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    if line_color is not None:
        tb.line.color.rgb = line_color
        tb.line.width = Pt(0.5)
    elif fill is None:
        tb.line.fill.background()
    return tb


def _add_logo(slide):
    png = _ensure_png(LOGO_SVG, size=256)
    if not png:
        return
    slide.shapes.add_picture(png, Inches(0.25), Inches(0.18),
                             height=Inches(0.55))


def _add_header_bar(slide, title: str):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, 0, SLIDE_W, Inches(0.95))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = CISCO_DARK

    _add_text(slide, Inches(1.1), Inches(0.18), Inches(11.5), Inches(0.6),
              title, size=22, bold=True, align=PP_ALIGN.LEFT, color=WHITE)

    _add_logo(slide)


def _add_device(slide, cx, cy, dev_type, hostname: str, platform: str = ""):
    left = int(cx - ICON_SIZE / 2)
    top  = int(cy - ICON_SIZE / 2)

    png = _icon_png(dev_type)
    if png:
        slide.shapes.add_picture(png, left, top,
                                 width=ICON_SIZE, height=ICON_SIZE)
    else:
        # Fallback: rounded rectangle with type label inside.
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     left, top, ICON_SIZE, ICON_SIZE)
        shp.fill.solid()
        shp.fill.fore_color.rgb = DEVICE_FALLBACK_COLORS.get(
            dev_type, DEVICE_FALLBACK_COLORS["unknown"])
        shp.line.color.rgb = CISCO_DARK
        shp.line.width = Pt(1.0)
        tf = shp.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = TYPE_LABELS.get(dev_type, "")
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = WHITE

    label_left = int(cx - LABEL_W / 2)
    label_top  = int(cy + ICON_SIZE / 2 + Emu(5000))
    block_h    = Inches(0.55)

    tb = slide.shapes.add_textbox(label_left, label_top, LABEL_W, block_h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(20000)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tb.fill.background()
    tb.line.fill.background()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = hostname
    run.font.size = Pt(9)
    run.font.bold = False
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    if platform:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = platform
        r2.font.size = Pt(9)
        r2.font.bold = True
        r2.font.color.rgb = RGBColor(0x00, 0x00, 0x00)


def _add_link(slide, x1, y1, x2, y2, local_if: str, remote_if: str):
    line = slide.shapes.add_connector(1, int(x1), int(y1), int(x2), int(y2))
    line.line.width = Pt(1.25)
    line.line.color.rgb = RGBColor(0x55, 0x55, 0x55)

    dx = x2 - x1
    dy = y2 - y1

    label_w = Inches(0.85)
    label_h = Inches(0.24)

    def _place(frac: float, text: str):
        if not text:
            return
        px = x1 + dx * frac
        py = y1 + dy * frac
        tb = slide.shapes.add_textbox(
            int(px - label_w / 2), int(py - label_h / 2),
            label_w, label_h)
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = Emu(20000)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tb.fill.solid()
        tb.fill.fore_color.rgb = WHITE
        tb.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        tb.line.width = Pt(0.5)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(8)
        run.font.bold = False
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    _place(0.20, local_if)
    _place(0.80, remote_if)


def build_topology_slide(prs: Presentation, local_name: str,
                         local_type: str, neighbors: List[Neighbor]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                0, Inches(0.95), SLIDE_W, SLIDE_H - Inches(0.95))
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_LIGHT

    _add_header_bar(slide, f"CDP Topology  -  {local_name}")

    cx = SLIDE_W / 2
    cy = SLIDE_H / 2 + Inches(0.1)

    grouped: dict[str, dict] = {}
    for n in neighbors:
        key = n.short_name
        g = grouped.setdefault(key, {"neighbor": n, "links": []})
        g["links"].append((n.local_interface, n.remote_interface))

    nodes = list(grouped.values())
    count = max(len(nodes), 1)

    radius_x = Inches(5.0)
    radius_y = Inches(2.3)

    angles = [-math.pi / 2 + (2 * math.pi * i / count) for i in range(count)]

    for i, item in enumerate(nodes):
        links = item["links"]
        angle = angles[i]
        nx = cx + radius_x * math.cos(angle)
        ny = cy + radius_y * math.sin(angle)

        for idx, (lif, rif) in enumerate(links):
            offset = (idx - (len(links) - 1) / 2) * Inches(0.22)
            perp_angle = angle + math.pi / 2
            ox = offset * math.cos(perp_angle)
            oy = offset * math.sin(perp_angle)
            _add_link(slide,
                      cx + ox, cy + oy,
                      nx + ox, ny + oy,
                      lif, rif)

    _add_device(slide, cx, cy, local_type, local_name)

    for i, item in enumerate(nodes):
        n: Neighbor = item["neighbor"]
        angle = angles[i]
        nx = cx + radius_x * math.cos(angle)
        ny = cy + radius_y * math.sin(angle)
        _add_device(slide, nx, ny, n.device_type, n.short_name, n.platform)


def build_matrix_slide(prs: Presentation, local_name: str,
                       neighbors: List[Neighbor]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                0, Inches(0.95), SLIDE_W, SLIDE_H - Inches(0.95))
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_LIGHT

    _add_header_bar(slide, f"Connectivity Matrix  -  {local_name}")

    headers = ["Local Interface", "Neighbor Device", "Remote Interface",
               "Platform", "Type", "IPv4", "Version"]
    col_widths = [Inches(w) for w in (1.4, 2.3, 1.6, 1.8, 1.0, 1.6, 3.0)]

    rows = len(neighbors) + 1
    cols = len(headers)
    left, top = Inches(0.3), Inches(1.15)
    height = Inches(0.42) * rows
    if height > Inches(6.2):
        height = Inches(6.2)

    table_shape = slide.shapes.add_table(rows, cols, left, top,
                                         sum(col_widths, Emu(0)), height)
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = CISCO_BLUE

    for r, n in enumerate(neighbors, start=1):
        values = [
            n.local_interface,
            n.short_name,
            n.remote_interface,
            n.platform,
            n.device_type,
            n.primary_ipv4,
            n.version,
        ]
        for c, val in enumerate(values):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val or ""
            run.font.size = Pt(9)
            run.font.color.rgb = TEXT_DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_ALT if r % 2 == 0 else WHITE


def build_presentation(local_name: str, neighbors: List[Neighbor],
                       output: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    local_type = "l3switch"

    build_topology_slide(prs, local_name, local_type, neighbors)
    build_matrix_slide(prs, local_name, neighbors)

    prs.save(output)


# ---------------------------------------------------------------------------
# Public helpers
# ---------------------------------------------------------------------------

def _guess_local_name(text: str) -> str:
    m = re.search(r"^([A-Za-z0-9._\-]+)\s*[#>]\s*sh(ow)?\s+cdp", text, re.MULTILINE)
    if m:
        return m.group(1)
    return "local-device"


def generate_pptx_from_cdp(text: str,
                           output_path: str,
                           local_name: Optional[str] = None) -> dict:
    """Parse the supplied CDP detail output and write a PPTX file.

    Returns a dict suitable for serializing to JSON containing status,
    message, and metadata.
    """
    if not text or not text.strip():
        return {"status": "error", "message": "No CDP input provided."}

    # Reset per-call warning bucket.
    _ICON_WARNINGS.clear()

    neighbors = parse_cdp_output(text)
    if not neighbors:
        return {"status": "error",
                "message": "No CDP neighbor entries were found in the input."}

    name = (local_name or "").strip() or _guess_local_name(text)
    build_presentation(name, neighbors, output_path)

    backend = _RENDERER_CACHE.get("backend") or "none"
    return {
        "status": "success",
        "message": f"Generated PowerPoint with {len(neighbors)} neighbor(s).",
        "local_name": name,
        "neighbor_count": len(neighbors),
        "output_path": output_path,
        "icon_dir": ICON_DIR,
        "renderer": backend,
        "warnings": list(_ICON_WARNINGS),
    }


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

SENTINEL = "END"


def _read_pasted_input() -> str:
    print("Paste the output of `show cdp neighbors detail` below.")
    print(f"When you are done, press Ctrl-D (macOS/Linux) or type '{SENTINEL}' "
          "on a new line and press Enter.\n")

    lines: List[str] = []
    while True:
        try:
            line = input()
        except EOFError:
            break
        if line.strip() == SENTINEL:
            break
        lines.append(line)
    return "\n".join(lines)


def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("input", nargs="?", default=None)
    ap.add_argument("-o", "--output", default="cdp_topology.pptx")
    ap.add_argument("-n", "--name", default=None)
    args = ap.parse_args(argv)

    if args.input is None:
        text = _read_pasted_input()
    elif args.input == "-":
        text = sys.stdin.read()
    else:
        with open(args.input, "r", encoding="utf-8", errors="replace") as fh:
            text = fh.read()

    result = generate_pptx_from_cdp(text, args.output, local_name=args.name)
    if result["status"] != "success":
        sys.exit(result["message"])
    print(f"\nWrote {args.output} ({result['neighbor_count']} neighbor entries, "
          f"local device: {result['local_name']})")


if __name__ == "__main__":
    main()
